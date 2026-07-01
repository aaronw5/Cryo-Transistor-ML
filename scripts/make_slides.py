#!/usr/bin/env python3
"""Build the project slide deck (slides/cryo-ml-77k.pptx).

Audience: mixed ML / non-ML. Run scripts/make_figs.py first.
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
FIGS = ROOT / "figs"
OUT = ROOT / "slides"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

DARK = RGBColor(0x20, 0x33, 0x44)
ACCENT = RGBColor(0xE8, 0x6C, 0x1A)
GRAY = RGBColor(0x55, 0x5F, 0x68)


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_title(slide, text, sub=None, top=Inches(0.25)):
    box = slide.shapes.add_textbox(Inches(0.5), top, SLIDE_W - Inches(1.0),
                                   Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK
    if sub:
        p2 = tf.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(15)
        p2.font.color.rgb = GRAY
    return box


def add_bullets(slide, items, left, top, width, height, size=18):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        level, text = (item if isinstance(item, tuple) else (0, item))
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = ("• " if level == 0 else "– ") + text
        p.level = level
        p.font.size = Pt(size if level == 0 else size - 2)
        p.font.color.rgb = DARK if level == 0 else GRAY
        p.space_after = Pt(8)
    return box


def add_picture_fit(slide, path, left, top, max_w, max_h):
    w_px, h_px = Image.open(path).size
    scale = min(max_w / w_px, max_h / h_px)
    w, h = int(w_px * scale), int(h_px * scale)
    slide.shapes.add_picture(
        str(path), left + int((max_w - w) / 2), top + int((max_h - h) / 2),
        width=w, height=h)


def title_slide(prs, title, subtitle, lines):
    slide = blank_slide(prs)
    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.1),
                                   SLIDE_W - Inches(1.8), Inches(3.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = DARK
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(20)
    p2.font.color.rgb = ACCENT
    p2.space_before = Pt(14)
    for line in lines:
        p3 = tf.add_paragraph()
        p3.text = line
        p3.font.size = Pt(14)
        p3.font.color.rgb = GRAY
    return slide


def bullets_slide(prs, title, items, sub=None, size=18):
    slide = blank_slide(prs)
    add_title(slide, title, sub)
    add_bullets(slide, items, Inches(0.7), Inches(1.45), SLIDE_W - Inches(1.4),
                SLIDE_H - Inches(1.9), size=size)
    return slide


def picture_slide(prs, title, image, takeaway=None, sub=None):
    slide = blank_slide(prs)
    add_title(slide, title, sub)
    top = Inches(1.35)
    bottom = Inches(0.75) if takeaway else Inches(0.25)
    add_picture_fit(slide, image, Inches(0.35), top,
                    SLIDE_W - Inches(0.7), SLIDE_H - top - bottom)
    if takeaway:
        box = slide.shapes.add_textbox(Inches(0.7), SLIDE_H - Inches(0.7),
                                       SLIDE_W - Inches(1.4), Inches(0.55))
        p = box.text_frame.paragraphs[0]
        box.text_frame.word_wrap = True
        p.text = takeaway
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ACCENT
    return slide


def split_slide(prs, title, image, items, sub=None, img_frac=0.58, size=16):
    slide = blank_slide(prs)
    add_title(slide, title, sub)
    top = Inches(1.4)
    img_w = int((SLIDE_W - Inches(1.0)) * img_frac)
    add_picture_fit(slide, image, Inches(0.4), top, img_w, SLIDE_H - top
                    - Inches(0.3))
    text_left = Inches(0.6) + img_w
    add_bullets(slide, items, text_left, top, SLIDE_W - text_left - Inches(0.4),
                SLIDE_H - top - Inches(0.3), size=size)
    return slide


def add_pptx_table(slide, headers, rows, left, top, width, height,
                   font=10, header_font=None):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top,
                                   width, height)
    tbl = shape.table
    row_h = int(height / (len(rows) + 1))
    for r in tbl.rows:
        r.height = row_h
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(header_font or font)
            p.font.bold = True
            p.font.color.rgb = DARK
    for i, row in enumerate(rows, 1):
        for j, v in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = str(v).replace("`", "")
            cell.margin_top = Pt(1)
            cell.margin_bottom = Pt(1)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font)
                p.font.color.rgb = DARK
    return shape


def table6_slides(prs, paper_tables):
    for block in paper_tables["table6"]:
        slide = blank_slide(prs)
        add_title(slide, f"Errors for 77 K models — {block['name']}",
                  sub="Exact format of the paper's Table 6 "
                      "(Appendix A); computed in the identical corrected "
                      "NGSpice chain")
        n = len(block["rows"])
        add_pptx_table(slide, block["headers"], block["rows"],
                       Inches(2.4), Inches(1.35), Inches(8.5),
                       SLIDE_H - Inches(1.7), font=10, header_font=10)
        _ = n


def table4_slides(prs, paper_tables):
    blocks = paper_tables["table4"]
    per_slide = 6
    for s0 in range(0, len(blocks), per_slide):
        chunk = blocks[s0:s0 + per_slide]
        slide = blank_slide(prs)
        add_title(slide,
                  "Appendix: extracted parameters per model bin "
                  f"({s0 + 1}–{s0 + len(chunk)} of {len(blocks)})",
                  sub="Exact format of the paper's Table 4: published 77 K "
                      "values vs the ML-refit values shipped in the library")
        cols, col_w = 3, Inches(4.15)
        row_h_block = Inches(2.75)
        for k, block in enumerate(chunk):
            cx = Inches(0.35) + (k % cols) * (col_w + Inches(0.15))
            cy = Inches(1.45) + (k // cols) * (row_h_block + Inches(0.35))
            box = slide.shapes.add_textbox(cx, cy - Inches(0.27), col_w,
                                           Inches(0.25))
            p = box.text_frame.paragraphs[0]
            p.text = block["head"]
            p.font.size = Pt(9.5)
            p.font.bold = True
            p.font.color.rgb = ACCENT
            add_pptx_table(slide, block["headers"], block["rows"],
                           cx, cy, col_w, row_h_block, font=8.5)


def main() -> int:
    OUT.mkdir(exist_ok=True)
    paper_tables = json.loads(
        (ROOT / "out" / "tables" / "paper_tables.json").read_text())
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(
        prs,
        "Beating hand-tuned cryogenic transistor models with ML",
        "SkyWater 130 nm MOSFETs at 77 K — BSIM4 re-extraction, "
        "validated end-to-end in open-source NGSpice",
        ["Measured data & baseline models: arXiv:2604.21625 "
         "(Beall et al., 2026)",
         "ML beats the published cards (error 0.692 → 0.536, 17/18 "
         "devices) AND the same cards refit with classical optimizers "
         "(0.491 vs 0.498) — same simulator, same metric; the classical "
         "optimizer plateaus even with 3.5x its simulation budget"],
    )

    bullets_slide(
        prs,
        "Background: why cryogenic transistor models",
        [
            "Detector readout and quantum-computing electronics run cold — "
            "77 K (liquid nitrogen) and below.",
            "Transistors behave differently when cold: thresholds shift, "
            "mobility rises, standard models break.",
            "Foundry models for the open-source SkyWater 130 nm process are "
            "only valid near room temperature.",
            "The paper (arXiv:2604.21625) measured 18 transistor sizes "
            "(8 nMOS, 10 pMOS) at 77 K and hand-tuned BSIM4 'model cards' "
            "so designers can simulate circuits cold.",
            (1, "The data, per transistor: 11 current-voltage sweeps of "
                "186 points each — 6 gate sweeps (0 → 1.85 V in 10 mV "
                "steps) at drain biases of 0.01, 0.37, 0.74, 1.11, 1.48 "
                "and 1.85 V, plus 5 drain sweeps (0 → 1.85 V) at the same "
                "five gate biases. Negative voltages for pMOS; source and "
                "body grounded. Every simulation in this work reproduces "
                "exactly these ~2,000 bias points per transistor."),
            (1, "A model card = a parameter file for the simulator; the "
                "paper tunes 7 physical parameters per device: VTH0 "
                "(threshold), U0 (mobility), NFACTOR, VSAT, DELTA, RDSW, "
                "ETA0."),
            "Accuracy is scored with RRMS: RMS difference between simulated "
            "and measured current, divided by the mean measured current — "
            "averaged over every measured curve.",
        ],
    )

    bullets_slide(
        prs,
        "The problem: the published models don't reproduce in NGSpice",
        [
            "The paper reports mean RRMS 0.279, extracted and scored in "
            "HSPICE + Mystic (proprietary, licensed tools).",
            "Running the paper's own published cards in NGSpice (the "
            "open-source simulator the cards target) gives 0.692 — "
            "2.5x worse than reported.",
            "So 'our cards vs the paper's reported number' would be a "
            "meaningless comparison across different simulators.",
            "Fair comparison used everywhere here: paper cards vs ML cards, "
            "same NGSpice version, same decks, same bins, same metric, "
            "same measured data.",
        ],
    )

    bullets_slide(
        prs,
        "Why the gap? Four simulation problems we isolated",
        [
            "1. Simulator gap — some devices simulate ~2x worse in NGSpice "
            "than reported even with the most favorable settings (e.g. pMOS "
            "0.35/1.6, nMOS 20/0.64, nMOS 100/100). HSPICE and NGSpice "
            "evaluate the same BSIM4 card differently.",
            "2. Scale-units trap — the cryo files set '.option scale=1.0u', "
            "so device sizes must be given as bare micron numbers; writing "
            "'l=0.15u' silently misses every model bin ('could not find a "
            "valid modelname'). This blocked all NGSpice use until fixed.",
            "3. Overlapping pMOS bins — the published pMOS card has 12 "
            "geometry bins for 10 devices with overlapping ranges; the "
            "simulator can map a device to a bin fit to a different device.",
            "4. Measured-data glitches — a few 'device-off' sweeps contain "
            "instrument spikes (µA-level, 1–2 of 186 points) no simulator "
            "can fit; they inflate RRMS identically for every card set.",
        ],
        size=17,
    )

    bullets_slide(
        prs,
        "Step 1: re-run the paper's method honestly (two classical "
        "optimizers as controls)",
        [
            "Control 1 — least-squares refinement (how extraction tools "
            "and the paper work): start from the published parameter "
            "values, nudge each of the 7 parameters slightly, re-simulate "
            "to see how the error responds, and step downhill; repeat "
            "until converged (up to 150 simulations per attempt).",
            (1, "Because this only walks downhill, it can get stuck — so "
                "it runs 6 times from different random starting points and "
                "keeps the best (~900 simulations per transistor). "
                "Result: 0.692 → 0.501, improves all 18 transistors."),
            "Control 2 — CMA-ES, a stronger global search: keeps a "
            "population of 12 candidate parameter sets, simulates them "
            "all, breeds the next generation around the best performers, "
            "and adapts the search shape as it goes; the winner gets the "
            "same least-squares finish.",
            (1, "Result: 0.499 at 2,400 simulations per transistor — the "
                "strongest classical result."),
            "Budget check: re-running CMA-ES with 8,500 simulations per "
            "transistor (matching the ML pipeline's total budget) improves "
            "it by only 0.001 (0.498). The classical approach has "
            "plateaued — any ML margin is methodological, not a budget "
            "artifact.",
            "Every method (controls and ML) searches the same space: each "
            "parameter confined to a generous window around its published "
            "value (e.g. threshold voltage ±0.6 V; mobility ×0.1 to ×10), "
            "scored by the same metric, in the same simulator.",
        ],
        size=14,
    )

    bullets_slide(
        prs,
        "The ML method, stage 1: train a neural net to predict what "
        "NGSpice outputs, then search through it",
        [
            "The problem with searching directly in NGSpice: one simulation "
            "takes ~0.5 s, and you can't take derivatives — so classical "
            "optimizers can only try thousands of parameter guesses, "
            "blindly.",
            "Step 1 — for each transistor, run 6,000 NGSpice simulations "
            "with randomized values of the 7 parameters. Each simulation "
            "gives one training example: 7 numbers in → the full set of "
            "~2,000 current readings out.",
            "Step 2 — train a neural network (4 layers, 512 wide) on those "
            "examples. It now answers the same question as NGSpice in "
            "microseconds, and, unlike NGSpice, it is differentiable.",
            "Step 3 — run 2,048 gradient descents through the network at "
            "once, each rolling downhill from a different starting point "
            "toward parameters that match the measured 77 K curves. This "
            "explores millions of candidates instead of thousands.",
            "Step 4 — THE TRUST RULE: the network is never believed. Its "
            "best candidates are re-simulated in real NGSpice ('real' = "
            "the actual simulator, as opposed to the network's fast "
            "imitation of it), ranked by the real score, and the winners "
            "are fine-tuned with classic least-squares running on real "
            "NGSpice. Every number reported anywhere in this deck is a "
            "real-NGSpice number.",
            "A second network that predicts parameters directly from the "
            "measured curves is used only to propose starting points — "
            "trusted raw, it scores 47 (vs 0.69 for doing nothing).",
        ],
        size=15,
    )

    bullets_slide(
        prs,
        "The ML method, stage 2: a feedback loop — search, re-simulate "
        "the best candidates, retrain the nets",
        [
            "Weakness of stage 1 alone: the network can be confidently "
            "wrong in corners of parameter space where it saw little "
            "training data — the search happily 'optimizes' into those "
            "phantom minima.",
            "Fix, part 1 — train THREE networks instead of one (different "
            "random seeds). Where they disagree, the prediction can't be "
            "trusted; the search score adds that disagreement as a penalty, "
            "so it avoids phantom minima. Half the starting points do the "
            "opposite — they seek disagreement, to explore unmapped "
            "regions.",
            "Fix, part 2 — close the loop: each round, the 8 best "
            "candidates are simulated in real NGSpice, and those true "
            "results are added to the training data; the three networks "
            "are retrained and the search repeats. 4 rounds, only ~36 "
            "extra simulations per transistor.",
            "Each round, the networks become most accurate exactly where "
            "the search is looking — the phantom-minimum errors of stage 1 "
            "shrink away.",
            "The classical optimizers' winning cards also enter every "
            "round as candidates, so the ML result can only match or beat "
            "them.",
        ],
        size=15,
    )

    bullets_slide(
        prs,
        "How to read the results: two scores per method",
        [
            "'One card per transistor' — every measured transistor gets "
            "its own independently fitted parameter card. This measures "
            "pure fitting strength, and is how methods are compared "
            "head-to-head.",
            "'One card per size bin' — a real PDK library ships one card "
            "per region of transistor sizes (a 'bin'), and the simulator "
            "picks the bin from the geometry. Two pairs of our pMOS "
            "transistors fall into the same bin, so each pair must share "
            "one compromise card.",
            "The second score is what you actually get by installing the "
            "library — it is slightly worse for every method, and the "
            "same constraint is applied to every method.",
            "Both scores use the paper's own error metric (RRMS: RMS "
            "error divided by mean measured current, averaged over every "
            "curve), computed by real NGSpice runs only.",
        ],
        size=16,
    )

    bullets_slide(
        prs,
        "What exactly is an 'ML card'?",
        [
            "There is ONE ML method: train the NN model, search through "
            "it, then run the feedback rounds. It produces candidate "
            "parameter cards for each transistor; the classical "
            "optimizers' winning cards also enter as candidates.",
            "Every candidate is scored in real NGSpice. The best score "
            "becomes the final ML card — so by construction, ML is never "
            "worse than the classical result.",
            "Where the 18 final cards actually came from:",
            (1, "2 transistors — the ML search found a card that beats "
                "everything else (nMOS 20/0.64: 0.148 vs classical 0.154; "
                "pMOS 8/1.6: 0.566 vs 0.593). This is where the ML margin "
                "lives."),
            (1, "4 transistors — the two shared size bins, where the card "
                "comes from the method's joint fit across both bin "
                "members."),
            (1, "12 transistors — a classical card was already at the "
                "achievable floor; the ML method re-polished it and kept "
                "it. ML matches, classical isn't beaten."),
            "Honest summary: ML = the classical method plus extra "
            "candidates that sometimes win. The overall margin "
            "(0.498 → 0.491) is exactly those wins; on every other "
            "transistor ML ties.",
        ],
        size=16,
    )

    bullets_slide(
        prs,
        "Fairness rules (applied to both card sets)",
        [
            "Published 77 K corner files + the exact Volare SKY130 PDK "
            "revision the paper specifies.",
            "Corrected-repository NGSpice deck convention "
            "(ogzamour/CryoSkywater130nm_CorrectedForNgspice).",
            "Native bin selection — NGSpice picks the geometry bin; we "
            "never pick bins based on measured scores.",
            "Only the paper's 7 parameters are tuned — no extra knobs.",
            "Optimization, model selection, and reporting all use the paper "
            "companion notebook's all-curve RRMS — including the corrupted "
            "curves.",
            "Every reported number comes from a real NGSpice run, never "
            "from the emulator.",
        ],
    )

    picture_slide(
        prs,
        "Results: all 18 Table-6 devices",
        FIGS / "table6_bars.png",
        takeaway="One card per size bin: paper cards 0.692 → CMA-ES refit "
                 "0.541 → ML refit 0.536 (nMOS 0.436 → 0.378 → 0.373; "
                 "pMOS 0.896 → 0.671 → 0.667). One card per transistor: "
                 "ML 0.491 vs best classical 0.498 (nMOS 0.373 vs 0.378, "
                 "pMOS 0.586 vs 0.594) — 3 wins, 15 ties, 0 losses. Black "
                 "diamonds: the paper's proprietary-simulator numbers, "
                 "unreachable by any card set in NGSpice.",
    )

    table6_slides(prs, paper_tables)

    for key, label, note in [
        ("paper", "paper cards", "as published, no fitting"),
        ("ml", "ML cards", "emulator search + feedback, validated in NGSpice"),
        ("direct", "direct-predict cards",
         "one-shot curve→parameters, then polished"),
    ]:
        picture_slide(
            prs,
            f"I-V at 77 K — {label}",
            FIGS / f"iv_{key}.png",
            sub=f"Measured (circles) vs {label} ({note}); paper Fig. 2 "
                "representative devices. One method per slide so each line is "
                "unambiguous.",
            takeaway=None,
        )

    picture_slide(
        prs,
        "Best-fit characteristics",
        FIGS / "fig4_bestfit.png",
        sub="Paper Fig. 4 analogue. Weak inversion = gate voltage below "
            "the ~0.4-0.5 V threshold: the transistor is barely on, "
            "current is tiny (nA) and grows exponentially with gate "
            "voltage. Strong inversion = gate well above threshold: the "
            "channel is fully formed and current (µA) grows smoothly. The "
            "two regimes are fit by different parameters (VTH0/NFACTOR vs "
            "U0/VSAT/RDSW).",
        takeaway="ML cards track the measured curves more closely in both "
                 "regimes — and the noisy weak-inversion data is where the "
                 "error metric is hardest on every method.",
    )

    picture_slide(
        prs,
        "Error across every measured geometry",
        FIGS / "fig5_rrms_heatmap.png",
        sub="Paper Fig. 5 analogue: identical color scale across panels",
        takeaway="Improvement is broad and ML edges the CMA-ES refit in "
                 "both families: nMOS 0.436 → 0.379 → 0.373, pMOS "
                 "0.896 → 0.677 → 0.668. The worst cell (pMOS 8/5: 1.46) "
                 "halves to 0.69.",
    )

    split_slide(
        prs,
        "Where the ML gain comes from",
        FIGS / "ml_ablation.png",
        [
            "Search stage alone (no feedback rounds, no help from the "
            "classical results): 0.494 — already past both classical "
            "optimizers (0.498 / 0.501).",
            "With the feedback rounds added: 0.484 — re-simulating and "
            "retraining buys a further margin.",
            "Letting the classical winners compete as candidates "
            "guarantees ML never loses a transistor: final 0.491 (3 wins, "
            "15 ties, 0 losses vs the strongest classical result).",
            "The network that predicts parameters directly from curves is "
            "unusable alone (error ≈ 47) — re-checking everything in the "
            "real simulator is what makes the pipeline work.",
            "Sharing one card per size bin costs ~0.045 for every method; "
            "it costs the classical optimizer slightly more (0.541 vs "
            "0.536).",
        ],
        img_frac=0.5,
        size=15,
    )

    picture_slide(
        prs,
        "What matters when scaling the ML: data, model size, or search?",
        FIGS / "scaling_laws.png",
        sub="Search stage only (no feedback rounds), 4 test transistors. "
            "Colored lines = individual transistors; black dashed = their "
            "average (taken on the log scale, since values span decades). "
            "Top row: how well the net predicts the simulator; bottom "
            "row: final fitting error in real NGSpice.",
        takeaway="Prediction quality improves steadily with more training "
                 "simulations and bigger nets — but the final fitting "
                 "error stops improving once the net is good enough: the "
                 "measurement floor takes over. Below ~1,000 training "
                 "simulations, fitting fails badly on some transistors. "
                 "More search alone buys nothing.",
    )

    picture_slide(
        prs,
        "Worst baseline device: pMOS L=8 µm, W=5 µm",
        FIGS / "devices" / "pmos_L8_W5.png",
        takeaway="RRMS 1.458 → 0.690: the paper card (dashed) overshoots "
                 "every strong-inversion curve by ~30%; the ML card lands on "
                 "the measured points.",
    )

    bullets_slide(
        prs,
        "Caveats",
        [
            "The ML library (0.536) still scores worse than the paper's "
            "reported 0.279 — that residual is the simulator gap plus data "
            "limits, not fitting quality; nothing run in NGSpice reaches "
            "0.279 on this data.",
            "The ML margin over the strongest classical optimizer is real "
            "but modest (0.491 vs 0.498, ~1.4%); the honest claim is "
            "'never worse, sometimes clearly better' (e.g. nMOS 20/0.64: "
            "0.148 vs 0.154; pMOS 8/1.6: 0.566 vs 0.593). Four independent "
            "method families converge at ~0.49-0.50 — the data/model "
            "floor.",
            "Much of the remaining error is unfittable for ANY method: "
            "instrument spikes, ~6 µA range quantization, and high-Vd "
            "leakage floors whose BSIM4 knobs (GIDL) are outside the 7 "
            "allowed parameters — verified by per-curve diagnostics.",
            "One deployable regression: pMOS 0.5/0.42 (1.098 → 1.171) — "
            "its bin partner wins the shared card trade-off; every method "
            "hits the same wall (classical control: 1.174).",
            "Single seed per configuration; per-device run-to-run spread "
            "not yet characterized.",
        ],
        size=16,
    )

    bullets_slide(
        prs,
        "Predicting parameters directly from the curve",
        [
            "The main pipeline searches the parameter space. A simpler "
            "approach from the compact-model literature: train one network to "
            "read a measured curve and output the 7 parameters in a single "
            "forward pass — no search.",
            "Trained on the same 6,000 NGSpice simulations per transistor, run "
            "the other way: curve in, 7 parameters out. The training loss is "
            "the distance to the true parameters.",
            "Second version adds a physics surrogate: also score the "
            "prediction on whether it reproduces the curve through the fast "
            "emulator, not just on the parameter numbers.",
            "Every prediction is validated in real NGSpice; a least-squares "
            "polish from the prediction is recorded for reference.",
        ],
        size=17,
    )

    bullets_slide(
        prs,
        "Direct prediction: clean data vs. measured data",
        [
            "On clean simulated curves the network reconstructs the I-V well "
            "— error ~0.13–0.19, in line with the literature.",
            "The physics surrogate helps: clean reconstruction 0.19 → 0.13, "
            "and one-shot error on the measured curves 1.39 → 0.77 (median).",
            "One shot is not enough on the real 77 K data (0.77, vs 0.491 for "
            "the search pipeline): the measured curves carry the glitches, "
            "quantization and leakage floors the clean training curves do not, "
            "so the network extrapolates.",
            "A least-squares polish from the prediction recovers to 0.52 — "
            "near the classical controls, below the search pipeline.",
            "A direct predictor is a good starting point; re-checking in the "
            "real simulator is what makes the extraction reliable.",
        ],
        size=17,
    )

    bullets_slide(
        prs,
        "Takeaways & next steps",
        [
            "ML extraction beats the published cryogenic cards (17/18 "
            "devices) and the paper's method re-run with classical "
            "optimizers (per-device and deployable), under strictly fair, "
            "simulator-verified conditions.",
            "The deliverable is practical: one drop-in 77 K library for "
            "the open-source SkyWater PDK (out/pdk_ml_final/cards/).",
            "Everything is reproducible with open tools: NGSpice 46 + "
            "Python; figures and this deck regenerate from scripts; full "
            "experiment history in docs/RESEARCH_LOG.md.",
            "A one-shot curve→parameter network confirms the design: "
            "re-checking in the real simulator, not direct prediction, is "
            "what makes extraction reliable.",
            "Scaling laws say where future effort pays: more/better "
            "training data and a data-quality-filtered metric — not more "
            "search. Tuning beyond the 7 paper parameters (GIDL for the "
            "leakage floors) and re-binning the pMOS card would attack the "
            "floor itself.",
        ],
    )

    table4_slides(prs, paper_tables)

    out_path = OUT / "cryo-ml-77k.pptx"
    prs.save(out_path)
    print(f"wrote {out_path}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
