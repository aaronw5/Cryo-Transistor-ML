#!/usr/bin/env python3
"""Build a minimal one-plot-per-slide results deck."""

from __future__ import annotations

import argparse
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import make_slide_plots


ROOT = Path(__file__).resolve().parents[1]
FIGS = ROOT / "figs"
SLIDE_PLOTS = ROOT / "slides" / "plots"
DEFAULT_OUTPUT = ROOT / "slides" / "cryo_ml_simple_results.pptx"

SLIDES = [
    ("Main I-V comparison", SLIDE_PLOTS / "main_iv_comparison.png"),
    ("All 18 transistors", SLIDE_PLOTS / "main_rrms_comparison.png"),
    ("Direct MLP", FIGS / "iv_direct.png"),
    ("Surrogate + FD", FIGS / "iv_emu_fd.png"),
    ("Finite-difference improvement",
     SLIDE_PLOTS / "fd_improvement_comparison.png"),
    ("Training-data scaling", FIGS / "scaling_laws.png"),
]


def add_title(slide, text: str, *, top: float = 0.22,
              height: float = 0.5, size: int = 22) -> None:
    box = slide.shapes.add_textbox(
        Inches(0.35), Inches(top), Inches(12.63), Inches(height))
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.name = "Arial"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(32, 36, 40)


def add_fitted_picture(slide, path: Path) -> None:
    left = Inches(0.22)
    top = Inches(0.82)
    max_width = Inches(12.89)
    max_height = Inches(6.47)
    with Image.open(path) as image:
        ratio = image.width / image.height

    width = max_width
    height = width / ratio
    if height > max_height:
        height = max_height
        width = height * ratio

    x = left + (max_width - width) / 2
    y = top + (max_height - height) / 2
    slide.shapes.add_picture(str(path), x, y, width=width, height=height)


def build(output: Path) -> None:
    make_slide_plots.build()
    missing = [path for _, path in SLIDES if not path.is_file()]
    if missing:
        raise FileNotFoundError(
            "missing slide figures: " + ", ".join(map(str, missing)))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    title_slide = prs.slides.add_slide(blank)
    add_title(title_slide, "Cryogenic SKY130 parameter extraction",
              top=2.72, height=0.75, size=30)
    subtitle = title_slide.shapes.add_textbox(
        Inches(1.2), Inches(3.55), Inches(10.93), Inches(0.55))
    paragraph = subtitle.text_frame.paragraphs[0]
    paragraph.text = "Measured 77 K data, paper cards, and ML-extracted cards"
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.name = "Arial"
    paragraph.font.size = Pt(18)
    paragraph.font.color.rgb = RGBColor(80, 86, 92)

    for title, path in SLIDES:
        slide = prs.slides.add_slide(blank)
        add_title(slide, title)
        add_fitted_picture(slide, path)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    print(f"wrote {output} ({len(prs.slides)} slides)")


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    args = parser.parse_args()
    build(args.output.resolve())


if __name__ == "__main__":
    main()
